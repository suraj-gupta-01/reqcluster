"""
Generate the ReqCluster Honeywell-hackathon pitch deck as a real, editable .pptx.

Design language (matches the product UI): near-black background, a single teal
accent, restrained glassmorphism (frosted translucent panels with a hairline
light border), strong type hierarchy, native editable flowcharts and cards.
No purple, no glossy bevels, no AI-slop tells. Visible text uses hyphens only.

Run:  uv run python docs/build_deck.py
Out:  docs/ReqCluster_Honeywell_Hackathon.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette -----
BG_TOP = RGBColor(0x0B, 0x10, 0x0F)
BG_BOT = RGBColor(0x07, 0x0A, 0x0A)
INK = RGBColor(0xEC, 0xF2, 0xF1)      # near-white text
MUTED = RGBColor(0x9A, 0xA7, 0xA4)    # secondary text
FAINT = RGBColor(0x6B, 0x78, 0x76)    # tertiary
TEAL = RGBColor(0x2F, 0xBC, 0xAA)     # single accent
TEAL_DEEP = RGBColor(0x0D, 0x81, 0x75)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SKY = RGBColor(0x38, 0xBD, 0xF8)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers -----
def _alpha(color_format, opacity_pct):
    """Add an <a:alpha> child to a ColorFormat (opacity 0-100)."""
    try:
        clr = color_format._color._xClr
        for a in clr.findall(qn("a:alpha")):
            clr.remove(a)
        el = clr.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))})
        clr.append(el)
    except Exception:
        pass


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(slide, x, y, w, h, fill, fill_alpha=100, line=None, line_alpha=100,
         line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if fill_alpha < 100:
            _alpha(sp.fill.fore_color, fill_alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
        if line_alpha < 100:
            _alpha(sp.line.color, line_alpha)
    _no_shadow(sp)
    return sp


def glass(slide, x, y, w, h, radius=0.06, fill_alpha=7, line_alpha=15):
    """A frosted translucent panel: faint white fill + hairline light border."""
    return rect(slide, x, y, w, h, fill=WHITE, fill_alpha=fill_alpha,
                line=WHITE, line_alpha=line_alpha, line_w=1.0, radius=radius)


def txt(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of paragraphs; each paragraph is a dict or list of run-dicts.
    run-dict: {t, size, color, bold, font, align, spacing, space_after}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    for i, para in enumerate(lines):
        runs = para if isinstance(para, list) else [para]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        meta = runs[0]
        p.alignment = meta.get("align", PP_ALIGN.LEFT)
        if "spacing" in meta:
            p.line_spacing = meta["spacing"]
        if "space_after" in meta:
            p.space_after = Pt(meta["space_after"])
        if "space_before" in meta:
            p.space_before = Pt(meta["space_before"])
        for rd in runs:
            r = p.add_run()
            r.text = rd["t"]
            f = r.font
            f.size = Pt(rd.get("size", 14))
            f.bold = rd.get("bold", False)
            f.name = rd.get("font", FONT)
            f.color.rgb = rd.get("color", INK)
    return tb


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=1.75, alpha=70):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if alpha < 100:
        _alpha(conn.line.color, alpha)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    _no_shadow(conn)
    return conn


def slide_base(eyebrow=None, title=None, subtitle=None, page=None):
    s = prs.slides.add_slide(BLANK)
    bg = rect(s, -0.06, -0.06, W + 0.12, H + 0.12, fill=BG_TOP, shape=MSO_SHAPE.RECTANGLE)
    bg.fill.gradient()
    try:
        stops = bg.fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = BG_TOP
        stops[1].position = 1.0
        stops[1].color.rgb = BG_BOT
        bg.fill.gradient_angle = 70.0
    except Exception:
        pass
    bg.line.fill.background()
    _no_shadow(bg)
    # soft teal depth glow, top-right
    glow = rect(s, 8.4, -2.4, 7.2, 6.0, fill=TEAL, fill_alpha=8,
                shape=MSO_SHAPE.OVAL)
    glow2 = rect(s, -2.6, 4.2, 6.0, 5.0, fill=TEAL_DEEP, fill_alpha=6,
                 shape=MSO_SHAPE.OVAL)
    if eyebrow:
        txt(s, 0.92, 0.62, 11.0, 0.3,
            [{"t": eyebrow.upper(), "size": 11, "color": TEAL, "bold": True}])
    if title:
        txt(s, 0.9, 0.92, 11.5, 1.0,
            [{"t": title, "size": 30, "color": INK, "bold": True}])
    if subtitle:
        txt(s, 0.92, 1.62, 11.4, 0.5,
            [{"t": subtitle, "size": 13.5, "color": MUTED, "spacing": 1.15}])
    # footer
    txt(s, 0.92, 7.04, 8.0, 0.3,
        [[{"t": "ReqCluster", "size": 9.5, "color": TEAL, "bold": True},
          {"t": "   Honeywell Hackathon · DP5", "size": 9.5, "color": FAINT}]])
    if page is not None:
        txt(s, 11.7, 7.04, 0.9, 0.3,
            [{"t": str(page), "size": 9.5, "color": FAINT, "align": PP_ALIGN.RIGHT}])
    return s


def card(slide, x, y, w, h, title, body, accent=TEAL, kicker=None):
    glass(slide, x, y, w, h, radius=0.07)
    rect(slide, x + 0.28, y + 0.34, 0.5, 0.07, fill=accent, radius=0.5)  # accent underline
    cy = y + 0.5
    if kicker:
        txt(slide, x + 0.28, cy, w - 0.5, 0.3,
            [{"t": kicker.upper(), "size": 9.5, "color": accent, "bold": True}])
        cy += 0.32
    txt(slide, x + 0.28, cy, w - 0.5, 0.5,
        [{"t": title, "size": 15.5, "color": INK, "bold": True}])
    txt(slide, x + 0.28, cy + 0.46, w - 0.56, h - (cy + 0.46 - y) - 0.2,
        [{"t": body, "size": 11.5, "color": MUTED, "spacing": 1.16}])


def metric(slide, x, y, w, h, number, label, accent=TEAL):
    glass(slide, x, y, w, h, radius=0.08)
    txt(slide, x + 0.25, y + 0.34, w - 0.5, 0.8,
        [{"t": number, "size": 34, "color": accent, "bold": True, "font": FONT}])
    txt(slide, x + 0.27, y + h - 0.78, w - 0.5, 0.6,
        [{"t": label, "size": 11, "color": MUTED, "spacing": 1.12}])


def node(slide, x, y, w, h, label, accent=TEAL, sub=None):
    glass(slide, x, y, w, h, radius=0.16, fill_alpha=9, line_alpha=18)
    rect(slide, x, y + 0.14, 0.06, h - 0.28, fill=accent, radius=0.5)  # left accent rail
    anchor_lines = [{"t": label, "size": 12.5, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}]
    if sub:
        anchor_lines.append({"t": sub, "size": 9, "color": MUTED, "align": PP_ALIGN.CENTER, "space_before": 2})
    txt(slide, x + 0.12, y, w - 0.2, h, anchor_lines, anchor=MSO_ANCHOR.MIDDLE)


def chips(slide, x, y, items, accent=TEAL, max_w=11.5, gap=0.16, ch=0.42):
    """Render a flowing row of small pill chips. Returns the y after wrapping."""
    cx = x
    cy = y
    for it in items:
        wpx = 0.22 + 0.092 * len(it)
        if cx + wpx > x + max_w:
            cx = x
            cy += ch + 0.16
        rect(slide, cx, cy, wpx, ch, fill=WHITE, fill_alpha=6, line=WHITE,
             line_alpha=14, radius=0.5)
        txt(slide, cx, cy, wpx, ch,
            [{"t": it, "size": 10.5, "color": INK, "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
        cx += wpx + gap
    return cy + ch


# ============================================================ SLIDE 1: TITLE ==
s = slide_base()
# brand mark
rect(s, 0.92, 2.5, 0.62, 0.62, fill=TEAL, radius=0.28)
rect(s, 0.92, 2.5, 0.62, 0.62, fill=None, line=WHITE, line_alpha=22, radius=0.28)
txt(s, 0.92, 2.5, 0.62, 0.62, [{"t": "R", "size": 24, "color": RGBColor(0x06, 0x12, 0x10),
                                "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.75, 2.5, 9.0, 0.7, [[{"t": "Req", "size": 30, "color": INK, "bold": True},
                              {"t": "Cluster", "size": 30, "color": TEAL, "bold": True}]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.92, 3.45, 11.2, 1.1,
    [{"t": "Requirements intelligence for systems engineering", "size": 26,
      "color": INK, "bold": True, "spacing": 1.05}])
txt(s, 0.94, 4.5, 10.8, 0.8,
    [{"t": "Automatically group, map dependencies between, and explain thousands of "
           "engineering requirements - in minutes, not days.", "size": 14,
      "color": MUTED, "spacing": 1.2}])
# bottom strip
rect(s, 0.92, 5.7, 11.5, 0.02, fill=WHITE, fill_alpha=10, shape=MSO_SHAPE.RECTANGLE)
txt(s, 0.92, 5.95, 11.5, 0.4,
    [[{"t": "HONEYWELL HACKATHON", "size": 11, "color": TEAL, "bold": True},
      {"t": "    DP5 · Grouping Functionally Related Requirements", "size": 11, "color": MUTED}]])
chips(s, 0.92, 6.45, ["SBERT", "UMAP", "HDBSCAN", "c-TF-IDF", "Dependency Graph",
                      "On-prem LLM", "MBSE Export"], max_w=11.5)

# ========================================================== SLIDE 2: PROBLEM ==
s = slide_base(eyebrow="The problem", page=2,
               title="Requirements pile up faster than humans can organize them",
               subtitle="A modern aerospace or automotive program carries 500 to 50,000 requirements, written by many teams over years.")
card(s, 0.92, 2.5, 3.72, 2.7, "Latency",
     "Manual grouping blocks downstream design, ICD production, and verification planning. It is the slow first step before real work begins.",
     accent=TEAL, kicker="Cost 1")
card(s, 4.8, 2.5, 3.72, 2.7, "Subjectivity",
     "Two analysts produce materially different cluster structures from the same set, so organizational consistency erodes review to review.",
     accent=AMBER, kicker="Cost 2")
card(s, 8.68, 2.5, 3.72, 2.7, "Incompleteness",
     "Hidden cross-cutting dependencies are missed by hand. Indirect relationships only surface through systematic similarity analysis.",
     accent=SKY, kicker="Cost 3")
metric(s, 0.92, 5.45, 3.72, 1.35, "8-12 hrs", "to cluster 500 requirements by hand", accent=ROSE)
metric(s, 4.8, 5.45, 3.72, 1.35, "~60%", "analyst agreement (Cohen's kappa)", accent=AMBER)
metric(s, 8.68, 5.45, 3.72, 1.35, "50,000", "requirements in a large program", accent=TEAL)

# ========================================================= SLIDE 3: SOLUTION ==
s = slide_base(eyebrow="Our solution", page=3,
               title="An end-to-end, explainable clustering pipeline",
               subtitle="Upload a requirements file and ReqCluster understands meaning, groups by function, and visualizes the result - fully automated.")
flow = [("Ingest", "CSV / XLSX"), ("SBERT", "384-d meaning"), ("UMAP", "reduce dims"),
        ("HDBSCAN", "density clusters"), ("c-TF-IDF", "auto labels"), ("Graph", "similarity")]
nx0, ny0, nw, nh = 0.92, 2.7, 1.78, 1.15
gap = (11.5 - len(flow) * nw) / (len(flow) - 1)
accents = [TEAL, TEAL, SKY, TEAL, AMBER, GREEN]
xs = []
for i, (lab, sub) in enumerate(flow):
    x = nx0 + i * (nw + gap)
    xs.append(x)
    node(s, x, ny0, nw, nh, lab, accent=accents[i], sub=sub)
    if i > 0:
        arrow(s, xs[i - 1] + nw, ny0 + nh / 2, x, ny0 + nh / 2)
txt(s, 0.92, 4.15, 11.5, 0.4,
    [{"t": "All deterministic. No black box. No API keys required - it runs fully offline.",
      "size": 12.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
card(s, 0.92, 4.8, 3.72, 1.95, "Interactive dashboard",
     "Scatter, similarity graph, searchable tables, and per-cluster detail - explore 50k points in the browser.", accent=TEAL)
card(s, 4.8, 4.8, 3.72, 1.95, "Human in the loop",
     "Analysts approve, correct, and reassign. Every correction teaches the next clustering pass.", accent=SKY)
card(s, 8.68, 4.8, 3.72, 1.95, "Standards export",
     "One click to ReqIF, SysML XMI, and Jama - results flow straight into existing RE toolchains.", accent=AMBER)

# ===================================================== SLIDE 4: DP5 DELIVERY ==
s = slide_base(eyebrow="DP5 deliverables", page=4,
               title="Exactly what the brief asks for - and more",
               subtitle="The Honeywell DP5 problem statement maps one-to-one onto what ReqCluster produces.")
card(s, 0.92, 2.5, 3.72, 3.1, "Logical grouping",
     "Requirements organized into clusters by functional similarity using SBERT embeddings and density-based clustering.",
     accent=TEAL, kicker="Deliverable 1")
card(s, 4.8, 2.5, 3.72, 3.1, "Dependency tree",
     "Hierarchical and sequential relationships inferred from each requirement's inputs, pre-conditions, and outputs - rendered as an interactive 2D/3D graph.",
     accent=SKY, kicker="Deliverable 2")
card(s, 8.68, 2.5, 3.72, 3.1, "Rationale document",
     "A generated explanation of why each requirement belongs to its group, plus a justification for every dependency edge.",
     accent=AMBER, kicker="Deliverable 3")
glass(s, 0.92, 5.85, 11.48, 0.95, radius=0.1)
txt(s, 1.2, 5.85, 11.0, 0.95,
    [[{"t": "On-prem LLM ready.  ", "size": 12.5, "color": TEAL, "bold": True},
      {"t": "Summaries and rationales run on a local open-source model (Qwen via Ollama) - no data leaves your network. Scales to 500 requirements across 150+ pages.",
       "size": 12.5, "color": MUTED}]], anchor=MSO_ANCHOR.MIDDLE)

# ==================================================== SLIDE 5: HOW IT WORKS ===
s = slide_base(eyebrow="How it works", page=5,
               title="From sentences to structure",
               subtitle="Each stage is a clean, testable step - the same proven recipe behind modern topic modeling.")
steps = [
    ("SBERT embeddings", "Every requirement becomes a 384-number vector. Similar meaning - similar vector.", TEAL),
    ("UMAP reduction", "384 dimensions shrink to 10 for clustering and 2 for the plot, keeping neighbors together.", SKY),
    ("HDBSCAN", "Finds dense groups automatically - no preset cluster count - and flags outliers as noise.", TEAL),
    ("c-TF-IDF labels", "Picks the words that make each cluster distinctive, so every group gets a clear name.", AMBER),
]
for i, (t, b, a) in enumerate(steps):
    x = 0.92 + (i % 2) * 5.84
    y = 2.6 + (i // 2) * 1.9
    card(s, x, y, 5.56, 1.7, t, b, accent=a, kicker=f"Step {i+1}")
txt(s, 0.92, 6.45, 11.5, 0.4,
    [{"t": "Cosine similarity then links related requirements into a network graph across cluster boundaries.",
      "size": 12, "color": MUTED, "align": PP_ALIGN.CENTER}])

# ==================================================== SLIDE 6: ARCHITECTURE ===
s = slide_base(eyebrow="Architecture", page=6,
               title="A clean three-tier system",
               subtitle="Separation of concerns keeps each layer independently testable and swappable.")
tiers = [
    ("Application", "React · Vite · Tailwind · Plotly (WebGL)", TEAL,
     "Upload · Overview · Scatter · Graph · Dependency Tree · Review Queue · Export"),
    ("API + ML", "FastAPI · SBERT · UMAP · HDBSCAN · c-TF-IDF · LLM services", SKY,
     "REST endpoints · async jobs · pipeline orchestration"),
    ("Data", "PostgreSQL / SQLite · Redis cache · embedding store", AMBER,
     "Sessions · requirements · clusters · graphs · constraints"),
]
ty = 2.55
for i, (t, sub, a, detail) in enumerate(tiers):
    y = ty + i * 1.42
    glass(s, 0.92, y, 11.48, 1.2, radius=0.08)
    rect(s, 0.92, y + 0.16, 0.07, 0.88, fill=a, radius=0.5)
    txt(s, 1.25, y + 0.18, 4.6, 0.9,
        [{"t": t, "size": 16, "color": INK, "bold": True},
         {"t": sub, "size": 10.5, "color": a, "space_before": 3}])
    txt(s, 6.0, y, 6.2, 1.2, [{"t": detail, "size": 11, "color": MUTED, "spacing": 1.12}],
        anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        arrow(s, 6.66, y + 1.2, 6.66, y + 1.42, color=FAINT, width=1.5, alpha=60)
txt(s, 0.92, 6.85, 11.5, 0.3,
    [{"t": "Offline-first · deterministic seeds · Docker + CI · 133 automated tests",
      "size": 11, "color": FAINT, "align": PP_ALIGN.CENTER}])

# ===================================================== SLIDE 7: HITL / AL =====
s = slide_base(eyebrow="Human in the loop + active learning", page=7,
               title="Experts correct it once, the model learns forever",
               subtitle="Corrections are not throwaway edits - they become constraints that reshape the next clustering pass.")
loop = ["Cluster", "Review & correct", "Constraints", "Re-cluster", "Quality improves"]
lx0, ly, lw, lh = 0.92, 2.75, 2.05, 1.05
lgap = (11.5 - len(loop) * lw) / (len(loop) - 1)
la = [TEAL, SKY, AMBER, TEAL, GREEN]
prevx = None
for i, lab in enumerate(loop):
    x = lx0 + i * (lw + lgap)
    node(s, x, ly, lw, lh, lab, accent=la[i])
    if prevx is not None:
        arrow(s, prevx + lw, ly + lh / 2, x, ly + lh / 2)
    prevx = x
card(s, 0.92, 4.3, 3.72, 2.4, "Must / cannot-link",
     "Each manual reassignment generates must-link and cannot-link pairs that the clusterer must honor.", accent=TEAL)
card(s, 4.8, 4.3, 3.72, 2.4, "Uncertainty sampling",
     "The least-confident assignments surface first, so review effort lands exactly where it matters.", accent=SKY)
card(s, 8.68, 4.3, 3.72, 2.4, "Quality tracking",
     "Silhouette and noise rate are recorded each iteration - you watch clustering get measurably better.", accent=GREEN)

# ======================================================= SLIDE 8: MBSE EXPORT =
s = slide_base(eyebrow="Integration", page=8,
               title="Results flow into the tools engineers already use",
               subtitle="Not a dead end - tool-importable exports in the standard requirements and model formats.")
ex = [
    ("ReqIF 1.2", "OMG XML with cluster-grouped hierarchy", "DOORS Next · ReqView · Polarion", TEAL),
    ("SysML / XMI", "UML 2.5 model, clusters as packages", "Papyrus · MagicDraw / Cameo", SKY),
    ("Jama Connect", "REST item + relationship bundle", "Jama (upload or offline JSON)", AMBER),
    ("CSV", "Flat table with cluster assignments", "Excel · Sheets · quick review", GREEN),
]
for i, (t, b, dest, a) in enumerate(ex):
    x = 0.92 + (i % 2) * 5.84
    y = 2.6 + (i // 2) * 2.0
    glass(s, x, y, 5.56, 1.78, radius=0.08)
    rect(s, x + 0.3, y + 0.32, 0.5, 0.07, fill=a, radius=0.5)
    txt(s, x + 0.3, y + 0.46, 5.0, 0.4, [{"t": t, "size": 15, "color": INK, "bold": True}])
    txt(s, x + 0.3, y + 0.92, 5.0, 0.4, [{"t": b, "size": 11, "color": MUTED}])
    txt(s, x + 0.3, y + 1.3, 5.0, 0.4,
        [[{"t": "Imports into  ", "size": 10, "color": FAINT},
          {"t": dest, "size": 10, "color": a, "bold": True}]])

# ====================================================== SLIDE 9: SCALE ========
s = slide_base(eyebrow="Built to scale", page=9,
               title="Engineered for 35k-50k+ requirements",
               subtitle="The architecture targets real program sizes, not just a demo dataset.")
metric(s, 0.92, 2.55, 2.74, 1.75, "50k", "requirements per run", accent=TEAL)
metric(s, 3.86, 2.55, 2.74, 1.75, "minutes", "on CPU, end to end", accent=SKY)
metric(s, 6.8, 2.55, 2.74, 1.75, "< 90s", "on GPU (cuML auto-detect)", accent=GREEN)
metric(s, 9.74, 2.55, 2.66, 1.75, "O(N log N)", "ANN similarity graph", accent=AMBER)
card(s, 0.92, 4.55, 3.72, 2.2, "Parallel + adaptive",
     "Multi-core UMAP and HDBSCAN; deterministic for small sets, fast-approximate for large ones; GPU when present.", accent=TEAL)
card(s, 4.8, 4.55, 3.72, 2.2, "Production datastore",
     "PostgreSQL with connection pooling and indexed hot paths; Redis-backed incremental embedding cache.", accent=SKY)
card(s, 8.68, 4.55, 3.72, 2.2, "Async jobs",
     "Clustering runs as a background job with live progress polling, so the UI never blocks on long runs.", accent=AMBER)

# ====================================================== SLIDE 10: WHY WIN =====
s = slide_base(eyebrow="Why ReqCluster", page=10,
               title="What the alternatives can't do",
               subtitle="Existing tools each solve a slice. ReqCluster is the end-to-end, explainable, supervised pipeline.")
comp = [
    ("DOORS / Jama", "Store and version requirements - but offer no intelligent clustering.", ROSE),
    ("k-means / LDA", "Need a preset cluster count and ignore domain-specific meaning.", AMBER),
    ("ChatGPT prompts", "Non-deterministic, manual, and disconnected from RE workflows.", SKY),
]
for i, (t, b, a) in enumerate(comp):
    card(s, 0.92 + i * 3.88, 2.5, 3.68, 2.05, t, b, accent=a)
glass(s, 0.92, 4.75, 11.48, 1.95, radius=0.09)
rect(s, 0.92, 4.95, 0.07, 1.55, fill=TEAL, radius=0.5)
txt(s, 1.3, 4.95, 11.0, 0.5, [{"t": "ReqCluster", "size": 17, "color": TEAL, "bold": True}])
txt(s, 1.3, 5.5, 10.9, 1.1,
    [{"t": "Domain-aware embeddings, automatic cluster discovery, deterministic explainable labels, "
           "a human-supervised correction loop, and tool-importable MBSE export - in one offline platform.",
      "size": 13, "color": INK, "spacing": 1.22}])

# ====================================================== SLIDE 11: IMPACT ======
s = slide_base(eyebrow="Impact", page=11,
               title="Measurable outcomes",
               subtitle="From days of manual effort to a few automated minutes, with results you can defend.")
metric(s, 0.92, 2.6, 3.72, 1.85, "99%", "less time (8-12 hrs to under 5 min)", accent=TEAL)
metric(s, 4.8, 2.6, 3.72, 1.85, "kappa - 1.0", "deterministic, repeatable grouping", accent=SKY)
metric(s, 8.68, 2.6, 3.72, 1.85, "133", "automated tests, green in CI", accent=GREEN)
metric(s, 0.92, 4.7, 3.72, 1.85, "0", "API keys - runs fully offline", accent=AMBER)
metric(s, 4.8, 4.7, 3.72, 1.85, "1-click", "Docker compose deployment", accent=TEAL)
metric(s, 8.68, 4.7, 3.72, 1.85, "4", "MBSE / RE export formats", accent=SKY)

# ====================================================== SLIDE 12: CLOSING =====
s = slide_base(page=12)
txt(s, 0.92, 2.35, 11.4, 0.4, [{"t": "THANK YOU", "size": 12, "color": TEAL, "bold": True}])
txt(s, 0.9, 2.75, 11.5, 1.2,
    [{"t": "ReqCluster turns a requirements backlog into a structured, explainable, "
           "engineer-ready map.", "size": 26, "color": INK, "bold": True, "spacing": 1.06}])
txt(s, 0.94, 4.2, 11.0, 0.5,
    [{"t": "Automated · explainable · human-supervised · offline-first · standards-ready.",
      "size": 14, "color": MUTED}])
glass(s, 0.92, 4.95, 11.48, 1.5, radius=0.09)
txt(s, 1.25, 5.1, 10.9, 0.4, [{"t": "BUILT WITH", "size": 10, "color": TEAL, "bold": True}])
chips(s, 1.25, 5.5, ["Python", "FastAPI", "sentence-transformers", "UMAP", "HDBSCAN",
                     "scikit-learn", "React", "Plotly", "PostgreSQL", "Redis", "Docker"],
      max_w=10.8)

out = os.path.join(os.path.dirname(__file__), "ReqCluster_Honeywell_Hackathon.pptx")
prs.save(out)
print(f"Saved {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
