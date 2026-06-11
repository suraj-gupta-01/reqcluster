"""
Build the ReqCluster Design-a-thon pitch deck (.pptx).

Premium dark-tech theme with a consistent design system:
gradient backgrounds, accent system, flowcharts (pipeline / 3-tier
architecture / feedback loop), stat cards and comparison layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
BG_TOP      = RGBColor(0x0A, 0x0F, 0x1F)   # near-black navy
BG_BOT      = RGBColor(0x13, 0x20, 0x47)   # deep blue
CARD        = RGBColor(0x16, 0x22, 0x40)   # card surface
CARD_HI     = RGBColor(0x1D, 0x2C, 0x52)   # lighter card
LINE        = RGBColor(0x2A, 0x3A, 0x63)   # hairline / border
WHITE       = RGBColor(0xF6, 0xF8, 0xFC)
MUTED       = RGBColor(0x9F, 0xB0, 0xCC)
FAINT       = RGBColor(0x6C, 0x7C, 0x9C)

CYAN        = RGBColor(0x35, 0xD6, 0xF0)   # primary accent
VIOLET      = RGBColor(0x8B, 0x7CF if False else 0x5C, 0xF6)
VIOLET      = RGBColor(0x8B, 0x5C, 0xF6)
GREEN       = RGBColor(0x35, 0xE0, 0x9B)
AMBER       = RGBColor(0xF6, 0xC1, 0x4B)
PINK        = RGBColor(0xF4, 0x6B, 0x9A)

ACCENTS = [CYAN, VIOLET, GREEN, AMBER, PINK]

FONT       = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_SEMI  = "Segoe UI Semibold"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def _noline(shape):
    shape.line.fill.background()


def shadow(shape, blur=120000, dist=55000, alpha=68000):
    """Soft outer shadow via raw XML (best-effort)."""
    try:
        spPr = shape._element.spPr
        for e in spPr.findall(qn('a:effectLst')):
            spPr.remove(e)
        eff = spPr.makeelement(qn('a:effectLst'), {})
        sh = eff.makeelement(qn('a:outerShdw'),
                             {'blurRad': str(blur), 'dist': str(dist),
                              'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '05070F'})
        a = clr.makeelement(qn('a:alpha'), {'val': str(alpha)})
        clr.append(a); sh.append(clr); eff.append(sh); spPr.append(eff)
    except Exception:
        pass


def rect(s, l, t, w, h, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=None, line_w=1.0, radius=0.08, shadowed=False):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        _noline(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadowed:
        shadow(sp)
    else:
        sp.shadow.inherit = False
    return sp


def grad(s, l, t, w, h, c1, c2, angle=90, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[1].position = 1.0
    stops[1].color.rgb = c2
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    _noline(sp)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font,tracking?)"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for run in para:
            txt, size, color, bold, font = run[0], run[1], run[2], run[3], run[4]
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def bg(s):
    grad(s, -0.06, -0.06, SW + 0.12, SH + 0.12, BG_TOP, BG_BOT, angle=120,
         shape=MSO_SHAPE.RECTANGLE)
    # decorative accent glows (soft rounded blocks, very subtle)
    blob = rect(s, SW - 3.0, -1.4, 4.2, 4.2, fill=RGBColor(0x16, 0x2A, 0x55),
                radius=0.5)
    blob.fill.fore_color.rgb = RGBColor(0x15, 0x27, 0x52)


def kicker(s, label, color=CYAN, l=0.9, t=0.62):
    # accent tick + eyebrow label
    rect(s, l, t + 0.02, 0.34, 0.10, fill=color, shape=MSO_SHAPE.RECTANGLE)
    text(s, l + 0.46, t - 0.12, 9, 0.4,
         [[(label.upper(), 12.5, color, True, FONT_SEMI)]])


def title(s, t1, t2=None, l=0.9, t=0.95, color2=CYAN):
    runs = [[(t1, 33, WHITE, True, FONT_SEMI)]]
    text(s, l, t, 11.5, 1.0, runs)
    if t2:
        text(s, l, t + 0.62, 11.5, 0.5,
             [[(t2, 15, MUTED, False, FONT)]])


def footer(s, n):
    rect(s, 0, SH - 0.04, SW, 0.06, fill=RGBColor(0x10, 0x1A, 0x33),
         shape=MSO_SHAPE.RECTANGLE)
    text(s, 0.9, SH - 0.52, 6, 0.3,
         [[("ReqCluster", 10, FAINT, True, FONT_SEMI),
           ("   ·  AI Requirements Clustering", 10, FAINT, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW - 2.4, SH - 0.52, 1.5, 0.3,
         [[(f"{n:02d} / 13", 10, FAINT, True, FONT_SEMI)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def chip(s, l, t, w, label, color=CYAN, txtcolor=None):
    h = 0.42
    c = rect(s, l, t, w, h, fill=CARD_HI, line=LINE, line_w=1.0, radius=0.5)
    text(s, l, t, w, h,
         [[(label, 12, txtcolor or WHITE, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return c


def badge_num(s, l, t, n, color):
    d = 0.5
    c = rect(s, l, t, d, d, fill=color, shape=MSO_SHAPE.OVAL)
    shadow(c, blur=70000, dist=25000, alpha=55000)
    text(s, l, t - 0.01, d, d, [[(str(n), 16, BG_TOP, True, FONT_SEMI)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- card grid
def feature_card(s, l, t, w, h, idx, head, body, accent):
    card = rect(s, l, t, w, h, fill=CARD, line=LINE, line_w=1.0,
                radius=0.07, shadowed=True)
    # top accent strip
    rect(s, l, t, w, 0.09, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    pad = 0.28
    if idx is not None:
        badge_num(s, l + pad, t + 0.30, idx, accent)
        head_t = t + 0.30
        head_l = l + pad + 0.66
        head_w = w - pad - 0.66 - 0.2
    else:
        head_t = t + 0.30
        head_l = l + pad
        head_w = w - 2 * pad
    text(s, head_l, head_t, head_w, 0.6,
         [[(head, 15.5, WHITE, True, FONT_SEMI)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, l + pad, t + 1.02, w - 2 * pad, h - 1.2,
         [[(b, 12.5, MUTED, False, FONT)] for b in body],
         line_spacing=1.05, space_after=5)


def bullet_block(s, l, t, w, items, accent):
    """Compact rows: accent square + text."""
    y = t
    for it in items:
        rect(s, l, y + 0.07, 0.13, 0.13, fill=accent,
             shape=MSO_SHAPE.RECTANGLE)
        text(s, l + 0.32, y - 0.04, w - 0.4, 0.5,
             [[(it, 13, RGBColor(0xCE, 0xD9, 0xEC), False, FONT)]],
             line_spacing=1.0)
        y += 0.52
    return y


# ================================================================ SLIDE 1
s = slide(); bg(s)
# big accent ring motif
grad(s, 8.6, 1.0, 5.4, 5.4, RGBColor(0x14, 0x2C, 0x5E), BG_BOT, angle=45,
     shape=MSO_SHAPE.OVAL, radius=0)
ring = s.shapes.add_shape(MSO_SHAPE.DONUT, Inches(8.9), Inches(1.3),
                          Inches(4.6), Inches(4.6))
ring.fill.solid(); ring.fill.fore_color.rgb = CYAN
try:
    ring.adjustments[0] = 0.06
except Exception:
    pass
_noline(ring); ring.shadow.inherit = False
ring2 = s.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.7), Inches(2.1),
                           Inches(3.0), Inches(3.0))
ring2.fill.solid(); ring2.fill.fore_color.rgb = VIOLET
try:
    ring2.adjustments[0] = 0.10
except Exception:
    pass
_noline(ring2); ring2.shadow.inherit = False
# small cluster dots inside ring
import math
cx, cy = 11.2, 3.6
for i, (ang, r, col) in enumerate([
        (20, 0.9, GREEN), (80, 1.1, CYAN), (150, 0.7, AMBER),
        (210, 1.0, PINK), (300, 0.85, WHITE), (340, 0.6, VIOLET)]):
    x = cx + r * math.cos(math.radians(ang))
    y = cy + r * math.sin(math.radians(ang))
    d = 0.22
    rect(s, x, y, d, d, fill=col, shape=MSO_SHAPE.OVAL)

kicker(s, "Design-a-thon  ·  AI Requirements Engineering", CYAN, l=0.9, t=1.5)
text(s, 0.86, 2.05, 8.2, 2.2,
     [[("Req", 70, WHITE, True, FONT_SEMI),
       ("Cluster", 70, CYAN, True, FONT_SEMI)]])
text(s, 0.9, 3.35, 8.0, 1.2,
     [[("AI-Powered Functional Requirement", 23, WHITE, False, FONT_LIGHT)],
      [("Clustering & Analysis Platform", 23, WHITE, False, FONT_LIGHT)]],
     line_spacing=1.05)
text(s, 0.9, 4.55, 7.6, 1.0,
     [[("Turn thousands of raw requirements into clear, traceable,",
        14.5, MUTED, False, FONT)],
      [("human-supervised structure — in minutes, not days.",
        14.5, MUTED, False, FONT)]],
     line_spacing=1.15)
# tech pills
pills = ["SBERT", "UMAP", "HDBSCAN", "c-TF-IDF", "Human-in-the-loop"]
x = 0.9
for i, p in enumerate(pills):
    w = 0.42 + len(p) * 0.115
    chip(s, x, 5.75, w, p, ACCENTS[i % len(ACCENTS)])
    x += w + 0.22
footer(s, 1)


# ================================================================ SLIDE 2
s = slide(); bg(s)
kicker(s, "The Problem")
title(s, "Requirements engineering does not scale",
      "As systems grow, manual organization becomes the bottleneck.")
stats = [("500–50k", "requirements per modern system, from many stakeholders", CYAN),
         ("25–40%", "of review-cycle time lost to manual clustering & dedup", AMBER),
         ("8–12 hrs", "for one analyst to cluster just 500 requirements", PINK)]
cw, gap = 3.74, 0.34
x = 0.9
for (num, lbl, col) in stats:
    card = rect(s, x, 2.0, cw, 1.85, fill=CARD, line=LINE, radius=0.08,
                shadowed=True)
    rect(s, x, 2.0, 0.09, 1.85, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.32, 2.18, cw - 0.5, 0.9,
         [[(num, 38, col, True, FONT_SEMI)]])
    text(s, x + 0.34, 3.0, cw - 0.6, 0.8,
         [[(lbl, 13, MUTED, False, FONT)]], line_spacing=1.05)
    x += cw + gap

text(s, 0.9, 4.25, 11.5, 0.4,
     [[("Three compounding pains", 15, WHITE, True, FONT_SEMI)]])
pains = [("Latency", "Requirement organization stalls downstream design, ICDs and verification planning.", CYAN),
         ("Subjectivity", "Different analysts produce materially different cluster structures from the same set.", VIOLET),
         ("Incompleteness", "Manual review misses the indirect semantic relationships automation finds reliably.", GREEN)]
x = 0.9
for (h, b, col) in pains:
    card = rect(s, x, 4.75, cw, 1.75, fill=CARD_HI, line=LINE, radius=0.08)
    rect(s, x + 0.3, 4.98, 0.34, 0.10, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.3, 5.12, cw - 0.6, 0.5, [[(h, 15, WHITE, True, FONT_SEMI)]])
    text(s, x + 0.3, 5.55, cw - 0.6, 0.9, [[(b, 12, MUTED, False, FONT)]],
         line_spacing=1.05)
    x += cw + gap
footer(s, 2)


# ================================================================ SLIDE 3
s = slide(); bg(s)
kicker(s, "Why Existing Tools Fall Short", PINK)
title(s, "The gap in today's toolchain",
      "Every category solves part of the problem — none solve it end to end.")
cards = [("Jama · IBM DOORS", ["Store and version requirements", "No intelligent clustering", "Traceability stays manual"], PINK),
         ("k-means · LDA", ["Ignore domain semantics", "Need a preset cluster count", "Fragile on real text"], AMBER),
         ("LLM chat tools", ["Non-deterministic output", "Manual prompt engineering", "Not integrated into RE flow"], VIOLET)]
cw, gap = 3.74, 0.34
x = 0.9
for (h, items, col) in cards:
    card = rect(s, x, 2.05, cw, 2.55, fill=CARD, line=LINE, radius=0.08,
                shadowed=True)
    rect(s, x, 2.05, cw, 0.09, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.3, 2.32, cw - 0.6, 0.5, [[(h, 16, WHITE, True, FONT_SEMI)]])
    text(s, x + 0.3, 2.9, cw - 0.6, 0.3,
         [[("MISSING", 10.5, col, True, FONT_SEMI)]])
    bullet_block(s, x + 0.3, 3.25, cw - 0.5, items, col)
    x += cw + gap
# ReqCluster answer banner
ban = grad(s, 0.9, 4.95, 11.52, 1.55, RGBColor(0x10, 0x3A, 0x4A), CARD_HI,
           angle=0)
ban.line.color.rgb = CYAN; ban.line.width = Pt(1.25)
text(s, 1.3, 5.18, 1.9, 1.1,
     [[("Req", 26, WHITE, True, FONT_SEMI), ("Cluster", 26, CYAN, True, FONT_SEMI)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, 3.5, 5.18, 8.6, 1.1,
     [[("The only pipeline that combines domain-aware embeddings, density-based",
        14, WHITE, True, FONT_SEMI)],
      [("clustering, and deterministic labeling — fully automated, fully explainable.",
        14, MUTED, False, FONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, 3)


# ================================================================ SLIDE 4
s = slide(); bg(s)
kicker(s, "Our Solution", GREEN)
title(s, "ReqCluster: intelligent, explainable, supervised",
      "Upload requirements — get organized, labeled, traceable clusters.")
feats = [(1, "Ingest & auto-group", ["Upload CSV or XLSX", "Auto-discovers functionally", "related requirement groups"], CYAN),
         (2, "Cluster with ML", ["SBERT embeddings +", "UMAP projection +", "HDBSCAN density clustering"], VIOLET),
         (3, "Explainable labels", ["Deterministic c-TF-IDF", "keyword labels on every", "cluster — no black box"], GREEN),
         (4, "Surface dependencies", ["Cosine similarity graph", "reveals hidden links", "between requirements"], AMBER)]
cw, gap = 2.78, 0.27
x = 0.9
for (i, h, b, col) in feats:
    feature_card(s, x, 2.05, cw, 2.5, i, h, b, col)
    x += cw + gap
# mission banner
mb = grad(s, 0.9, 4.85, 11.52, 1.5, RGBColor(0x12, 0x2B, 0x55),
          RGBColor(0x16, 0x22, 0x40), angle=0)
mb.line.color.rgb = GREEN; mb.line.width = Pt(1.25)
text(s, 1.25, 4.98, 2.4, 1.25,
     [[("MISSION", 12, GREEN, True, FONT_SEMI)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 3.0, 4.98, 9.1, 1.25,
     [[("Cut requirement review-cycle time by at least ", 16, WHITE, False, FONT),
       ("60%", 16, GREEN, True, FONT_SEMI),
       (" — while keeping a human in command.", 16, WHITE, False, FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)


# ================================================================ SLIDE 5  (pipeline flowchart)
s = slide(); bg(s)
kicker(s, "How It Works  ·  Core ML Pipeline", CYAN)
title(s, "From raw text to labeled, connected clusters",
      "A deterministic, end-to-end pipeline — fixed seed (42) for reproducible runs.")
stages = [("SBERT", "Embeddings", "384-dim vectors\ncontent-hash cached", CYAN),
          ("UMAP", "Reduction", "→ 10-D cluster\n→ 2-D visualize", VIOLET),
          ("HDBSCAN", "Clustering", "dense clusters\n+ noise, no preset k", GREEN),
          ("c-TF-IDF", "Labeling", "top keywords →\nauto cluster label", AMBER),
          ("Graph", "Similarity", "cosine > threshold\nhidden links", PINK)]
n = len(stages)
nw, nh = 2.05, 2.0
gapx = (11.52 - nw * n) / (n - 1)
top = 2.5
x = 0.9
for i, (t1, t2, sub, col) in enumerate(stages):
    card = rect(s, x, top, nw, nh, fill=CARD, line=LINE, radius=0.10,
                shadowed=True)
    rect(s, x, top, nw, 0.55, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(s, x, top + 0.30, nw, 0.28, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x, top + 0.05, nw, 0.5, [[(t1, 16, BG_TOP, True, FONT_SEMI)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, top + 0.68, nw, 0.4, [[(t2, 14, WHITE, True, FONT_SEMI)]],
         align=PP_ALIGN.CENTER)
    text(s, x + 0.12, top + 1.12, nw - 0.24, 0.8,
         [[(line, 11, MUTED, False, FONT)] for line in sub.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=1)
    # arrow to next
    if i < n - 1:
        ax = x + nw + (gapx - 0.34) / 2
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(ax),
                                Inches(top + nh/2 - 0.22), Inches(0.34),
                                Inches(0.44))
        ar.fill.solid(); ar.fill.fore_color.rgb = LINE
        _noline(ar); ar.shadow.inherit = False
    x += nw + gapx
# input + output captions
text(s, 0.9, top + nh + 0.28, 11.52, 0.5,
     [[("INPUT  ", 11.5, FAINT, True, FONT_SEMI),
       ("requirement texts[]", 12.5, MUTED, False, FONT),
       ("           the heavy work runs in a threadpool — UI polls step-by-step progress           ",
        11.5, FAINT, False, FONT),
       ("OUTPUT  ", 11.5, FAINT, True, FONT_SEMI),
       ("clusters · assignments · graph", 12.5, MUTED, False, FONT)]],
     align=PP_ALIGN.CENTER)
footer(s, 5)


# ================================================================ SLIDE 6  (architecture tiers)
s = slide(); bg(s)
kicker(s, "System Architecture", VIOLET)
title(s, "A clean three-tier platform",
      "Separated concerns, offline-first — cloud or on-prem LLMs are optional.")
tiers = [("TIER 3", "Application",
          "React + Vite + Tailwind dashboard  ·  Plotly scatter  ·  similarity-graph UI  ·  sidebar workspace",
          CYAN),
         ("TIER 2", "API & ML",
          "FastAPI REST (/api)  ·  service layer  ·  ML core (embeddings · reduction · clustering · labeling · graph)  ·  LLM services",
          VIOLET),
         ("TIER 1", "Data",
          "CSV / XLSX preprocessing  ·  SQLite via SQLAlchemy  ·  embedding cache (.npy)",
          GREEN)]
top = 2.1
th = 1.24
gapy = 0.2
x = 0.9
w = 11.52
for i, (tg, name, body, col) in enumerate(tiers):
    y = top + i * (th + gapy)
    card = rect(s, x, y, w, th, fill=CARD, line=LINE, radius=0.07,
                shadowed=True)
    rect(s, x, y, 2.55, th, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.07)
    rect(s, x + 2.2, y, 0.4, th, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.3, y, 2.2, th,
         [[(tg, 13, BG_TOP, True, FONT_SEMI)], [(name, 20, BG_TOP, True, FONT_SEMI)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=2)
    text(s, x + 2.95, y, w - 3.25, th,
         [[(body, 13.5, RGBColor(0xCE, 0xD9, 0xEC), False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    # down arrow between tiers
    if i < len(tiers) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.5),
                                Inches(y + th - 0.02), Inches(0.3),
                                Inches(gapy + 0.04))
        ar.fill.solid(); ar.fill.fore_color.rgb = LINE
        _noline(ar); ar.shadow.inherit = False
text(s, 0.9, top + 3 * (th + gapy) + 0.12, 11.52, 0.4,
     [[("Each tier exposes clean interfaces — enabling independent development, testing and scaling.",
        12.5, MUTED, False, FONT)]], align=PP_ALIGN.CENTER)
footer(s, 6)


# ============================================ generic phase slide builder
def phase_slide(n, kick, kcol, t1, t2, cards, note=None, note_col=CYAN):
    s = slide(); bg(s)
    kicker(s, kick, kcol)
    title(s, t1, t2)
    cw, gap = 2.78, 0.27
    x = 0.9
    ncards = len(cards)
    if ncards == 4:
        cw, gap = 2.78, 0.27
    elif ncards == 3:
        cw, gap = 3.74, 0.34
    top = 2.15
    ch = 2.75 if not note else 2.5
    for i, (h, b, col) in enumerate(cards):
        feature_card(s, x, top, cw, ch, i + 1, h, b, col)
        x += cw + gap
    if note:
        nb = grad(s, 0.9, top + ch + 0.25, 11.52, 1.15,
                  RGBColor(0x12, 0x2B, 0x55), RGBColor(0x16, 0x22, 0x40), angle=0)
        nb.line.color.rgb = note_col; nb.line.width = Pt(1.0)
        text(s, 1.25, top + ch + 0.25, 10.9, 1.15,
             [[(note, 14, WHITE, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n)
    return s


# ================================================================ SLIDE 7
phase_slide(
    7, "Phase 2  ·  LLM Semantic Enrichment", CYAN,
    "Smarter embeddings — optional and safe",
    "Expand each requirement with domain context, then cluster on richer meaning.",
    [("Provider choice", ["Mock (offline default),", "local Ollama / Qwen, or", "OpenAI-compatible gateway"], CYAN),
     ("3 embedding modes", ["Cluster on base,", "enriched, or hybrid", "embeddings"], VIOLET),
     ("Bounded & cached", ["Strict JSON parsing,", "bounded retries,", "content-hash caching"], GREEN),
     ("Measured lift", ["Quality scoring +", "embedding comparison", "and ablation metrics"], AMBER)],
    note="Any provider failure falls back to the deterministic offline path — the pipeline never breaks.",
    note_col=CYAN)

# ================================================================ SLIDE 8
phase_slide(
    8, "Phase 3  ·  Automated Cluster Refinement", VIOLET,
    "ClusterLLM-style merge & split suggestions",
    "The system proposes structural improvements — the analyst stays the decision-maker.",
    [("Merge candidates", ["Centroid similarity", "+ silhouette delta", "identify over-split groups"], CYAN),
     ("Split candidates", ["Spread + GMM", "bimodality + silhouette", "find blended clusters"], VIOLET),
     ("Representatives", ["Surfaces points closest", "to each centroid for", "fast human review"], GREEN),
     ("Full audit trail", ["Accept / reject, with a", "before-and-after", "refinement audit log"], AMBER)],
    note="Every suggestion is explainable, and every decision is logged for traceability.",
    note_col=VIOLET)


# ================================================================ SLIDE 9 (feedback loop)
s = slide(); bg(s)
kicker(s, "Phase 4  ·  Human-in-the-Loop", GREEN)
title(s, "The analyst stays in control",
      "Corrections become constraints that reshape the model — a closing loop.")
loop = [("Correct", "Move requirements;\nadjust cluster sizes", CYAN),
        ("Constrain", "Must-link / cannot-link\nML constraints", VIOLET),
        ("Re-cluster", "Constrained run\nrebuilds + relabels", GREEN),
        ("Validate", "Conflict & cycle\nchecks; track quality", AMBER)]
nw, nh = 2.45, 1.7
top = 2.55
gapx = (11.52 - nw * 4) / 3
x = 0.9
centers = []
for i, (h, b, col) in enumerate(loop):
    card = rect(s, x, top, nw, nh, fill=CARD, line=LINE, radius=0.10,
                shadowed=True)
    rect(s, x, top, nw, 0.62, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(s, x, top + 0.34, nw, 0.28, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x, top + 0.06, nw, 0.55, [[(h, 16, BG_TOP, True, FONT_SEMI)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.12, top + 0.78, nw - 0.24, 0.85,
         [[(ln, 11.5, MUTED, False, FONT)] for ln in b.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=1)
    centers.append((x + nw / 2, x, x + nw))
    if i < 3:
        ax = x + nw + (gapx - 0.32) / 2
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(ax),
                                Inches(top + nh / 2 - 0.2), Inches(0.32),
                                Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = LINE
        _noline(ar); ar.shadow.inherit = False
    x += nw + gapx
# return loop arrow (bottom) from Validate back to Correct
ly = top + nh + 0.45
# vertical down from last card
rect(s, centers[3][0] - 0.015, top + nh, 0.03, ly - (top + nh),
     fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
# horizontal back
rect(s, centers[0][0], ly, centers[3][0] - centers[0][0], 0.03,
     fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
# up arrow into first
upa = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(centers[0][0] - 0.13),
                         Inches(top + nh + 0.02), Inches(0.26),
                         Inches(0.45))
upa.fill.solid(); upa.fill.fore_color.rgb = GREEN
_noline(upa); upa.shadow.inherit = False
text(s, 0.9, ly - 0.32, 11.52, 0.3,
     [[("Active-learning loop  ·  uncertainty sampling feeds the next review",
        12, GREEN, True, FONT_SEMI)]], align=PP_ALIGN.CENTER, wrap=False)
text(s, 0.9, ly + 0.18, 11.52, 0.4,
     [[("Rejection cleanly rolls back the change. Human judgment is encoded into the model, not discarded.",
        12.5, MUTED, False, FONT)]], align=PP_ALIGN.CENTER)
footer(s, 9)


# ================================================================ SLIDE 10
phase_slide(
    10, "Phase 5 + DP5  ·  Dependencies · Active Learning · MBSE", AMBER,
    "Closing the loop — and connecting to system design",
    "From clusters to a traceable, design-ready model of the whole requirement set.",
    [("Dependency tree", ["Infers hierarchical,", "sequential & data deps", "+ rationale document"], CYAN),
     ("Active learning", ["Re-clusters under", "constraints; uncertainty-", "samples the next queue"], VIOLET),
     ("Quality history", ["Silhouette, noise rate &", "cluster count tracked", "across every iteration"], GREEN),
     ("MBSE export", ["ReqIF · SysML/UML XMI ·", "Jama · CSV — into DOORS,", "Papyrus, MagicDraw, Jama"], AMBER)],
    note="Standards-based export turns clustered requirements into artifacts your existing engineering tools can import.",
    note_col=AMBER)


# ================================================================ SLIDE 11  (impact)
s = slide(); bg(s)
kicker(s, "Impact & Results", GREEN)
title(s, "Measurable outcomes vs. industry baseline",
      "Faster, more consistent, more complete — and verified.")
cards = [("8–12 hrs", "< 5 min", "to cluster 500 requirements", CYAN),
         ("~60%", "> 90%", "cross-analyst consistency (deterministic)", VIOLET),
         ("manual", "automatic", "hidden-relationship detection via graph", GREEN),
         ("130+", "tests · CI", "validated on every push to the repo", AMBER)]
cw, gap = 2.78, 0.27
x = 0.9
top = 2.15
ch = 2.7
for (base, now, lbl, col) in cards:
    card = rect(s, x, top, cw, ch, fill=CARD, line=LINE, radius=0.08,
                shadowed=True)
    rect(s, x, top, cw, 0.09, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.28, top + 0.34, cw - 0.5, 0.4,
         [[("BASELINE", 10, FAINT, True, FONT_SEMI)]])
    text(s, x + 0.28, top + 0.62, cw - 0.5, 0.5,
         [[(base, 18, MUTED, True, FONT_SEMI)]])
    ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x + 0.3),
                            Inches(top + 1.18), Inches(0.26), Inches(0.32))
    ar.fill.solid(); ar.fill.fore_color.rgb = col
    _noline(ar); ar.shadow.inherit = False
    text(s, x + 0.28, top + 1.55, cw - 0.5, 0.6,
         [[(now, 30, col, True, FONT_SEMI)]])
    text(s, x + 0.28, top + 2.12, cw - 0.5, 0.55,
         [[(lbl, 11.5, MUTED, False, FONT)]], line_spacing=1.0)
    x += cw + gap
# bottom strip
bb = grad(s, 0.9, top + ch + 0.25, 11.52, 1.05, RGBColor(0x10, 0x33, 0x3F),
          RGBColor(0x16, 0x22, 0x40), angle=0)
bb.line.color.rgb = GREEN; bb.line.width = Pt(1.0)
text(s, 1.25, top + ch + 0.25, 10.9, 1.05,
     [[("Explainable by construction", 14, GREEN, True, FONT_SEMI),
       ("  —  deterministic seeds, c-TF-IDF labels and full audit logs make every result reproducible and reviewable.",
        13.5, WHITE, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11)


# ================================================================ SLIDE 12  (tech + principles)
s = slide(); bg(s)
kicker(s, "Tech Stack & Design Principles", CYAN)
title(s, "Engineered to be trustworthy",
      "Open-source, deterministic, and built to scale.")
# left: stack groups
groups = [("Backend & API", ["Python 3.11", "FastAPI", "Uvicorn"], CYAN),
          ("Machine Learning", ["sentence-transformers", "UMAP", "HDBSCAN", "scikit-learn"], VIOLET),
          ("Frontend", ["React", "Vite", "Tailwind CSS", "Plotly.js"], GREEN),
          ("Data", ["SQLite", "SQLAlchemy", "embedding cache"], AMBER)]
gx = 0.9
gy = 2.05
gw = 5.6
text(s, gx, gy - 0.05, gw, 0.4, [[("STACK", 12, FAINT, True, FONT_SEMI)]])
yy = gy + 0.4
for (h, items, col) in groups:
    text(s, gx, yy, gw, 0.32, [[(h, 13.5, WHITE, True, FONT_SEMI)]])
    cx = gx
    cy = yy + 0.36
    for it in items:
        w = 0.42 + len(it) * 0.108
        if cx + w > gx + gw:
            cx = gx; cy += 0.5
        chip(s, cx, cy, w, it, col)
        cx += w + 0.16
    yy = cy + 0.62

# right: principles
px = 6.95
pw = 5.45
text(s, px, gy, pw, 0.4, [[("DESIGN PRINCIPLES", 12, FAINT, True, FONT_SEMI)]])
prins = [("Explainability", "c-TF-IDF labels, no black-box output", CYAN),
         ("Human oversight", "validation & correction built in", VIOLET),
         ("Open-source first", "no vendor lock-in", GREEN),
         ("Determinism", "fixed seeds (42), reproducible", AMBER),
         ("Incremental delivery", "value from Phase 1 onward", PINK),
         ("Scalability", "batch embed, cache, edge limits", CYAN)]
yy = gy + 0.45
for (h, b, col) in prins:
    row = rect(s, px, yy, pw, 0.62, fill=CARD, line=LINE, radius=0.18)
    rect(s, px + 0.22, yy + 0.21, 0.18, 0.18, fill=col, shape=MSO_SHAPE.OVAL)
    text(s, px + 0.6, yy, pw - 0.8, 0.62,
         [[(h + "   ", 13, WHITE, True, FONT_SEMI),
           (b, 11.5, MUTED, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.72
footer(s, 12)


# ================================================================ SLIDE 13  (closing)
s = slide(); bg(s)
# big ring motif again, left
ring = s.shapes.add_shape(MSO_SHAPE.DONUT, Inches(-1.4), Inches(4.0),
                          Inches(4.2), Inches(4.2))
ring.fill.solid(); ring.fill.fore_color.rgb = RGBColor(0x16, 0x2A, 0x55)
try:
    ring.adjustments[0] = 0.10
except Exception:
    pass
_noline(ring); ring.shadow.inherit = False
kicker(s, "Vision & Closing", CYAN, l=0.9, t=1.5)
text(s, 0.86, 2.0, 11.6, 2.0,
     [[("From raw requirements to", 40, WHITE, True, FONT_LIGHT)],
      [("traceable, design-ready ", 40, WHITE, True, FONT_LIGHT),
       ("system models", 40, CYAN, True, FONT_SEMI)]],
     line_spacing=1.05)
text(s, 0.9, 4.0, 10.8, 1.0,
     [[("Five incremental phases, each shipping standalone value — explainable AI",
        15.5, MUTED, False, FONT)],
      [("paired with human supervision, the right model for safety-critical engineering.",
        15.5, MUTED, False, FONT)]],
     line_spacing=1.2)
# thank you + cta
tb = grad(s, 0.9, 5.45, 6.0, 1.05, RGBColor(0x12, 0x2B, 0x55),
          RGBColor(0x10, 0x33, 0x3F), angle=0)
tb.line.color.rgb = CYAN; tb.line.width = Pt(1.0)
text(s, 0.9, 5.45, 6.0, 1.05,
     [[("Thank you", 22, WHITE, True, FONT_SEMI),
       ("    ·    questions welcome", 15, MUTED, False, FONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 7.2, 5.45, 5.2, 1.05,
     [[("ReqCluster", 16, CYAN, True, FONT_SEMI)],
      [("AI Requirements Clustering & Analysis", 12, MUTED, False, FONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
footer(s, 13)


# ---------------------------------------------------------------- save
out = r"c:\Workspace\reqcluster\ReqCluster_Designathon.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
